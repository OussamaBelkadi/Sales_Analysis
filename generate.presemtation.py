import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import os
from datetime import datetime

# --- Load Data ---
file_path = "sales_data.xlsx"
df = pd.read_excel(file_path)

# --- Preprocessing ---
df.columns = df.columns.str.strip()
df = df.rename(columns={
    "OrderDate": "Order Date",
    "Ship Date": "Ship Date",
    "Order Quantity": "Order Quantity",
    "Unit Selling Price": "Unit Selling Price",
    "Unit Cost": "Unit Cost",
})

# Convert dates
df["Order Date"] = pd.to_datetime(df["Order Date"])
df["Year"] = df["Order Date"].dt.year
df["Month"] = df["Order Date"].dt.to_period("M")

# Calculated columns
df["Sales"] = df["Order Quantity"] * df["Unit Selling Price"]
df["Cost"] = df["Order Quantity"] * df["Unit Cost"]
df["Profit"] = df["Sales"] - df["Cost"]
df["Profit Margin"] = np.where(df["Sales"] == 0, 0, df["Profit"] / df["Sales"]) * 100

# --- Date Table ---
date_range = pd.date_range(df["Order Date"].min(), df["Order Date"].max())
date_table = pd.DataFrame({"Date": date_range})
date_table["Year"] = date_table["Date"].dt.year
date_table["Month"] = date_table["Date"].dt.to_period("M")

# --- Output folder ---
os.makedirs("output", exist_ok=True)

# Setup style
sns.set_theme(style="whitegrid")

# --- 1) Ventes par produit et comparaison avec l'année précédente ---
sales_by_product = df.groupby(["Product Description Index", "Year"])["Sales"].sum().unstack()
sales_by_product.plot(kind="bar", figsize=(14, 7), colormap="viridis")
plt.title("Ventes par produit - comparaison annuelle")
plt.ylabel("Montant des ventes")
plt.xlabel("Produit")
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig("output/ventes_par_produit_comparaison.png")
plt.close()

# --- Ajouter un graphique en rond (camembert) pour les produits ---
total_sales_by_product = df.groupby("Product Description Index")["Sales"].sum().sort_values(ascending=False).head(5)
plt.figure(figsize=(8, 8))
plt.pie(total_sales_by_product, labels=total_sales_by_product.index, autopct='%1.1f%%', startangle=140, colors=sns.color_palette("pastel"))
plt.title("Répartition des ventes - Top 5 Produits")
plt.tight_layout()
plt.savefig("output/pie_top5_produits.png")
plt.close()

# --- 2) Ventes par mois et comparaison avec l'année précédente ---
sales_by_month = df.groupby(["Month", "Year"])["Sales"].sum().unstack()
sales_by_month.index = sales_by_month.index.astype(str)
sales_by_month.plot(figsize=(14, 7), marker='o')
plt.title("Ventes par mois - comparaison annuelle")
plt.ylabel("Montant des ventes")
plt.xlabel("Mois")
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig("output/ventes_par_mois_comparaison.png")
plt.close()

# --- 3) Ventes des 5 premières villes ---
city_sales = df.groupby("Delivery Region Index")["Sales"].sum().sort_values(ascending=False).head(5)
city_sales.plot(kind="bar", figsize=(10, 6), color='skyblue')
plt.title("Top 5 des villes par ventes")
plt.ylabel("Montant des ventes")
plt.xlabel("Ville")
plt.tight_layout()
plt.savefig("output/top_5_villes.png")
plt.close()

# --- 4) Bénéfice par canal comparé à l'année précédente ---
profit_by_channel = df.groupby(["Channel", "Year"])["Profit"].sum().unstack()
profit_by_channel.plot(kind="bar", figsize=(12, 6), colormap="plasma")
plt.title("Bénéfice par canal - comparaison annuelle")
plt.ylabel("Montant du profit")
plt.xlabel("Canal")
plt.tight_layout()
plt.savefig("output/profit_par_canal.png")
plt.close()

# --- 5) Top 5 ventes par client ---
top_clients = df.groupby(["Customer Name Index", "Year"])["Sales"].sum().unstack()
top_clients = top_clients.sort_values(by=top_clients.columns[-1], ascending=False).head(5)
top_clients.plot(kind="bar", figsize=(12, 6), colormap="coolwarm")
plt.title("Top 5 clients par ventes - comparaison annuelle")
plt.ylabel("Montant des ventes")
plt.xlabel("Client")
plt.tight_layout()
plt.savefig("output/top_5_clients.png")
plt.close()

# --- 6) Last 5 ventes par client ---
last_clients = df.groupby(["Customer Name Index", "Year"])["Sales"].sum().unstack()
last_clients = last_clients.sort_values(by=last_clients.columns[-1], ascending=True).head(5)
last_clients.plot(kind="bar", figsize=(12, 6), colormap="autumn")
plt.title("Last 5 clients par ventes - comparaison annuelle")
plt.ylabel("Montant des ventes")
plt.xlabel("Client")
plt.tight_layout()
plt.savefig("output/last_5_clients.png")
plt.close()

# --- 7) Cartes de KPI ---
kpis = {
    "Total des ventes": df["Sales"].sum(),
    "Total des profits": df["Profit"].sum(),
    "Marge bénéficiaire (%)": df["Profit"].sum() / df["Sales"].sum() * 100,
    "Quantité totale vendue": df["Order Quantity"].sum()
}

plt.figure(figsize=(10, 6))
kpi_labels = list(kpis.keys())
kpi_values = list(kpis.values())
colors = sns.color_palette("Set2")
plt.barh(kpi_labels, kpi_values, color=colors)
plt.title("Cartes d'indicateurs de performance (KPI)")
plt.xlabel("Valeurs")
plt.tight_layout()
plt.savefig("output/kpi_cards.png")
plt.close()

# --- Create PowerPoint presentation ---
prs = Presentation()

# Function to add a slide with title and text
def add_slide_with_title(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.placeholders[1].text = content

# Function to add an image slide
def add_image_slide(prs, image_path):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), height=Inches(5.5))

# --- Title slide ---
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Analyse des ventes"
subtitle = slide.placeholders[1]
subtitle.text = "Projet d'analyse des ventes basé sur des données réelles"

# --- Ventes par produit ---
add_slide_with_title(prs, "Ventes par produit", "Voici les ventes par produit et leur comparaison avec l'année précédente.")
add_image_slide(prs, "output/ventes_par_produit_comparaison.png")

# --- Répartition des ventes par produits ---
add_slide_with_title(prs, "Répartition des ventes - Top 5 Produits", "Répartition des ventes parmi les 5 produits les plus populaires.")
add_image_slide(prs, "output/pie_top5_produits.png")

# --- Ventes par mois ---
add_slide_with_title(prs, "Ventes par mois", "Voici les tendances des ventes par mois, avec comparaison d'année en année.")
add_image_slide(prs, "output/ventes_par_mois_comparaison.png")

# --- Top 5 des villes ---
add_slide_with_title(prs, "Top 5 des villes par ventes", "Les ventes par ville, montrant les 5 principales villes.")
add_image_slide(prs, "output/top_5_villes.png")

# --- Bénéfice par canal ---
add_slide_with_title(prs, "Bénéfice par canal", "Analyse des bénéfices par canal de distribution.")
add_image_slide(prs, "output/profit_par_canal.png")

# --- Top 5 clients ---
add_slide_with_title(prs, "Top 5 clients", "Les 5 clients ayant généré le plus de ventes.")
add_image_slide(prs, "output/top_5_clients.png")

# --- Last 5 clients ---
add_slide_with_title(prs, "Last 5 clients", "Les 5 clients ayant généré le moins de ventes.")
add_image_slide(prs, "output/last_5_clients.png")

# --- KPI cards ---
add_slide_with_title(prs, "Indicateurs de performance (KPI)", "Voici les principaux indicateurs de performance liés aux ventes.")
add_image_slide(prs, "output/kpi_cards.png")

# --- Conclusion ---
add_slide_with_title(prs, "Conclusion", "Ce projet d'analyse des ventes a permis d'extraire des informations précieuses à partir des données brutes, en offrant des perspectives sur les performances de l'entreprise à travers différentes dimensions. Grâce à l'analyse des ventes par produit, par mois, et par canal de distribution, nous avons pu identifier les produits les plus performants, les tendances saisonnières des ventes, ainsi que les régions géographiques et les clients les plus rentables. L'ajout des indicateurs de performance clés (KPI) a également mis en lumière la rentabilité globale, la marge bénéficiaire, et l'évolution des volumes de ventes au fil du temps. Les résultats obtenus offrent une base solide pour prendre des décisions éclairées visant à optimiser les stratégies de vente, améliorer la gestion des stocks, et maximiser les profits futurs.")

# Save the presentation
prs.save("sales_analysis_presentation.pptx")

print("Présentation générée avec succès.")
